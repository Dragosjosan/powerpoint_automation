import os
import sys
import json
import traceback
from pathlib import Path
from typing import Dict, Any

from loguru import logger
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image

logger.add(
    sys.stderr,
    format="<green>{time:YYYY-MM-DD HH:mm:ss}</green> | <level>{level: <8}</level> | <level>{message}</level>",
    level="DEBUG",
    backtrace=True,  # Show traceback for errors
    diagnose=True,  # Show variables in traceback
    catch=True,  # Catch exceptions in the logging handler
)


class PowerPointTemplateError(Exception):
    pass


class SlideNotFoundError(PowerPointTemplateError):
    pass


class ImageNotFoundError(PowerPointTemplateError):
    pass


class TableNotFoundError(PowerPointTemplateError):
    pass


def load_data(data_path: str) -> Dict[str, Any]:
    try:
        if not Path(data_path).exists():
            logger.error(f"Data file not found: {data_path}")
            raise FileNotFoundError(f"Data file not found: {data_path}")
        with open(data_path) as f:
            return json.load(f)

    except Exception as e:
        logger.exception(f"Failed to load data from {data_path}")
        raise PowerPointTemplateError(f"Failed to load data: {e}") from e


class PowerPointTemplate:
    def __init__(self, template_path: str):
        self.template_path = self._initialize_template(template_path)
        self.prs = Presentation(self.template_path)
        self.slide_map = self._create_slide_map()

    def _initialize_template(self, template_path: str) -> str:
        if not Path(template_path).exists():
            logger.exception(f"Template file not found: {template_path}")
            raise FileNotFoundError(f"Template file not found: {template_path}")
        logger.debug(f"Template file found: {template_path}")
        return template_path

    def _create_slide_map(self) -> Dict[str, Any]:
        try:
            logger.debug("Creating slide map")
            slide_map = {}
            for i, slide in enumerate(self.prs.slides):
                if slide.shapes.title:
                    logger.debug(f"Slide {i} has title: {slide.shapes.title.text}")
                    title = slide.shapes.title.text
                    slide_map[title] = {"index": i, "slide": slide}
            logger.debug(f"Slide map created")
            return slide_map
        except Exception as e:
            logger.exception("Failed to create slide map")
            raise PowerPointTemplateError(f"Failed to create slide map: {e}") from e

    def apply_data(self, data_dict: Dict[str, Any]) -> None:
        try:
            for slide_name, slide_data in data_dict.items():
                logger.debug(f"Try to apply data to slide: '{slide_name}'")
                if slide_name not in self.slide_map:
                    logger.warning(f"Slide '{slide_name}' not found in template")
                    continue

                slide = self.slide_map[slide_name]["slide"]

                try:
                    if "text" in slide_data:
                        logger.debug(f"Input data contains text data for slide: '{slide_name}'")
                        self._replace_text_placeholders(slide, slide_data["text"])

                    if "tables" in slide_data:
                        logger.debug(f"Input data contains table data for slide: '{slide_name}'")
                        self._update_tables(slide, slide_data["tables"])

                    if "images" in slide_data:
                        logger.debug(f"Input data contains image data for slide: '{slide_name}'")
                        self._replace_images(slide, slide_data["images"])
                except Exception as e:
                    logger.exception(f"Error processing slide '{slide_name}'")
                    continue
        except Exception as e:
            logger.exception("Failed to apply data to presentation")
            raise PowerPointTemplateError(f"Failed to apply data: {e}") from e

    def _replace_text_placeholders(self, slide, text_data: Dict[str, str]) -> None:
        try:
            for shape in slide.shapes:
                if not hasattr(shape, "text"):
                    continue

                try:
                    text = shape.text
                    for key, value in text_data.items():
                        placeholder = "{{" + key + "}}"
                        if placeholder in text:
                            logger.debug(f"Replacing placeholder: '{placeholder}' with '{value}'")
                            text = text.replace(placeholder, str(value))

                    if text != shape.text:
                        shape.text = text
                except Exception as e:
                    logger.exception("Error processing shape in slide")
                    continue
        except Exception as e:
            logger.exception("Failed to replace text placeholders")
            raise PowerPointTemplateError(f"Failed to replace text: {e}") from e

    def _update_tables(self, slide, tables_data: Dict[str, Any]) -> None:
        try:
            table_shapes = [shape for shape in slide.shapes if shape.has_table]
            if not table_shapes:
                raise TableNotFoundError("No tables found in slide")

            logger.debug(f"Found {len(table_shapes)} tables in slide")

            for table_index, table_data in tables_data.items():
                try:
                    table_index = int(table_index)
                    table = None

                    if table_index < len(table_shapes):
                        table = table_shapes[table_index].table
                    else:
                        for shape in table_shapes:
                            if (hasattr(shape, "name") and shape.name == table_index) or (
                                table_data.get("identifier") and table_data["identifier"] in shape.text
                            ):
                                table = shape.table
                                break

                    if table and "data" in table_data:
                        rows = min(len(table_data["data"]), len(table.rows))
                        for r in range(rows):
                            row_data = table_data["data"][r]
                            cols = min(len(row_data), len(table.columns))
                            for c in range(cols):
                                table.cell(r, c).text = str(row_data[c])
                    else:
                        logger.warning(f"Table {table_index} not found or no data provided")
                except Exception as e:
                    logger.exception(f"Error updating table {table_index}")
                    continue
        except TableNotFoundError as e:
            logger.warning(str(e))
        except Exception as e:
            logger.exception("Failed to update tables")
            raise PowerPointTemplateError(f"Failed to update tables: {e}") from e

    def _replace_images(self, slide, images_data: Dict[str, str]) -> None:
        try:
            for image_name, image_path in images_data.items():
                try:
                    if not os.path.exists(image_path):
                        raise ImageNotFoundError(f"Image file '{image_path}' not found")

                    if image_name.isdigit():
                        image_index = int(image_name)
                        images = [shape for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE]
                        if image_index < len(images):
                            self._replace_single_image(slide, images[image_index], image_path)
                            continue
                        placeholders = [
                            shape for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER
                        ]
                        if image_index < len(placeholders):
                            self._replace_single_image(slide, placeholders[image_index], image_path)
                            continue
                        logger.warning(f"No suitable image placeholder found for index {image_index}")
                    else:
                        found = False
                        for shape in slide.shapes:
                            if (
                                shape.shape_type == MSO_SHAPE_TYPE.PICTURE
                                or shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER
                            ) and (
                                (hasattr(shape, "name") and shape.name == image_name)
                                or (hasattr(shape, "alt_text") and shape.alt_text == image_name)
                            ):
                                self._replace_single_image(slide, shape, image_path)
                                found = True
                                break
                        if not found:
                            logger.warning(f"No image placeholder found with name/alt_text: {image_name}")
                except ImageNotFoundError as e:
                    logger.warning(str(e))
                    continue
                except Exception as e:
                    logger.exception(f"Error replacing image {image_name}")
                    continue
        except Exception as e:
            logger.exception("Failed to replace images")
            raise PowerPointTemplateError(f"Failed to replace images: {e}") from e

    def _replace_single_image(self, slide, shape, image_path: str) -> None:
        """Replace a single image with error handling."""
        try:
            left, top, box_width, box_height = shape.left, shape.top, shape.width, shape.height

            sp = shape._element
            sp.getparent().remove(sp)

            with Image.open(image_path) as img:
                img_width, img_height = img.size

            width_ratio = box_width / img_width
            height_ratio = box_height / img_height
            scale = min(width_ratio, height_ratio)

            new_width = int(img_width * scale)
            new_height = int(img_height * scale)

            new_left = left + int((box_width - new_width) / 2)
            new_top = top + int((box_height - new_height) / 2)

            slide.shapes.add_picture(image_path, new_left, new_top, new_width, new_height)
        except Exception as e:
            logger.exception("Failed to replace single image")
            raise PowerPointTemplateError(f"Failed to replace single image: {e}") from e

    def save(self, output_path: str) -> None:
        """Save the presentation with error handling."""
        try:
            self.prs.save(output_path)
            logger.info(f"Presentation saved as {output_path}")
        except Exception as e:
            logger.exception("Failed to save presentation")
            raise PowerPointTemplateError(f"Failed to save presentation: {e}") from e


def main():
    try:
        logger.info("Starting PowerPoint template application")
        template_path = "example-presentation.pptx"
        template = PowerPointTemplate(template_path)

        logger.debug("Available slides in template:")
        for title in template.slide_map.keys():
            logger.info(f"- {title}")

        logger.info("Loading input data")
        data = load_data("data.json")

        logger.info("Applying data to template")
        template.apply_data(data)

        output_path = "updated_presentation.pptx"
        template.save(output_path)
        logger.info(" Presentation updated successfully!")

    except FileNotFoundError:
        logger.exception("File not found error")
        sys.exit(1)
    except PowerPointTemplateError:
        logger.exception("PowerPoint template error")
        sys.exit(1)
    except Exception:
        logger.exception("Unexpected error")
        sys.exit(1)


if __name__ == "__main__":
    main()
